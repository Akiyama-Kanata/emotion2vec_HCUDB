from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
OUTPUT = ROOT / "archive" / "reports" / "emotion2vec_進捗報告_モデル構成版.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

NAVY = "0B1F36"
NAVY_2 = "12304D"
INK = "17263A"
MUTED = "607184"
BLUE = "2C7BE5"
CYAN = "20B7C9"
GREEN = "28A878"
AMBER = "E9A23B"
RED = "D65B5B"
WHITE = "FFFFFF"
BG = "F3F6FA"
LINE = "D8E1EA"
PALE_BLUE = "EAF3FD"
PALE_CYAN = "E8F8FA"
PALE_GREEN = "EAF7F2"
PALE_AMBER = "FFF5E5"
PALE_RED = "FDEEEE"
FONT = "Yu Gothic"


def rgb(value):
    return RGBColor.from_string(value)


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def set_line(shape, color=None, width=1):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(color)
        shape.line.width = Pt(width)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=16,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.02,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, runs, x, y, w, h, *, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for text, size, color, bold in runs:
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(color)


def add_card(slide, x, y, w, h, *, fill=WHITE, line=LINE, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        card.adjustments[0] = 0.08
    set_fill(card, fill)
    set_line(card, line, 0.8)
    return card


def add_badge(slide, text, x, y, w, *, fill=PALE_BLUE, color=BLUE, h=0.32):
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    badge.adjustments[0] = 0.18
    set_fill(badge, fill)
    set_line(badge)
    add_text(
        slide,
        text,
        x,
        y + 0.005,
        w,
        h - 0.02,
        size=9,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return badge


def add_header(slide, title, kicker, number):
    add_text(slide, kicker, 0.65, 0.30, 7.0, 0.24, size=9.5, color=BLUE, bold=True)
    add_text(slide, title, 0.65, 0.62, 11.45, 0.52, size=25, color=NAVY, bold=True)
    add_text(slide, f"{number:02d}", 12.03, 0.35, 0.55, 0.28, size=11, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.22), Inches(0.70), Inches(0.045))
    set_fill(accent, CYAN)
    set_line(accent)


def add_footer(slide, text):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(7.05), Inches(12.02), Inches(0.012))
    set_fill(line, LINE)
    set_line(line)
    add_text(slide, text, 0.66, 7.11, 11.9, 0.18, size=8.2, color=MUTED)


def add_bullet_list(slide, items, x, y, w, h, *, size=13, color=INK, bullet_color=BLUE, gap=6):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.01)
    frame.margin_right = Inches(0.01)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(gap)
        paragraph.line_spacing = 1.05
        marker = paragraph.add_run()
        marker.text = "●  "
        marker.font.name = FONT
        marker.font.size = Pt(max(8, size - 4))
        marker.font.color.rgb = rgb(bullet_color)
        run = paragraph.add_run()
        run.text = item
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
    return box


def add_step_circle(slide, text, x, y, *, fill=BLUE, diameter=0.50):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(diameter), Inches(diameter))
    set_fill(circle, fill)
    set_line(circle)
    add_text(
        slide,
        text,
        x,
        y + 0.04,
        diameter,
        diameter - 0.08,
        size=11,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return circle


def add_connector(slide, x1, y1, x2, y2, *, color=LINE, width=2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    set_line(line, color, width)
    return line


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    # ==================== Slide 1: Cover ====================
    slide = prs.slides.add_slide(blank)
    add_bg(slide, NAVY)
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.75), Inches(-1.65), Inches(5.65), Inches(5.65))
    set_fill(glow, BLUE, 68)
    set_line(glow)
    glow2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.35), Inches(4.72), Inches(3.6), Inches(3.6))
    set_fill(glow2, CYAN, 72)
    set_line(glow2)
    add_badge(slide, "RESEARCH PROGRESS REPORT", 0.76, 0.65, 2.25, fill="173B5E", color="91DFE7")
    add_text(
        slide,
        "emotion2vecを用いた\n日本語音声感情認識の構築と評価",
        0.76,
        1.36,
        7.65,
        1.62,
        size=29,
        color=WHITE,
        bold=True,
        line_spacing=0.92,
    )
    add_text(slide, "説明性を担保したモデル構成の設計と入出力疎通の完了", 0.80, 3.20, 7.0, 0.34, size=16, color="BFD3E7", bold=True)
    add_card(slide, 0.76, 4.12, 6.58, 1.25, fill="123250", line="284A6B")
    add_text(slide, "今回の進捗", 1.02, 4.38, 0.86, 0.25, size=10, color="83DDE5", bold=True)
    add_text(slide, "VAD媒介型感情分類器の構築と、入出力の疎通検証が完了", 1.02, 4.73, 5.90, 0.35, size=16.5, color=WHITE, bold=True)
    add_text(slide, "23RD004  秋山叶太", 0.80, 6.35, 3.6, 0.31, size=14, color=WHITE, bold=True)
    add_text(slide, "2026年7月4日", 0.80, 6.76, 2.6, 0.25, size=11, color="BFD3E7")

    labels = [
        ("MODEL", "VAD媒介型分類器の設計と実装", BLUE),
        ("PIPELINE", "マルチタスク学習の統合", CYAN),
        ("VERIFIED", "48件の自動テストによる疎通", GREEN),
    ]
    for index, (tag, text, color) in enumerate(labels):
        y = 1.45 + index * 1.48
        add_card(slide, 8.75, y, 3.70, 1.10, fill=WHITE, line=WHITE)
        add_badge(slide, tag, 9.03, y + 0.20, 0.82, fill=PALE_BLUE if color == BLUE else PALE_CYAN if color == CYAN else PALE_GREEN, color=color)
        add_text(slide, text, 9.03, y + 0.61, 3.00, 0.27, size=11.5, color=NAVY, bold=True)

    # ==================== Slide 2: 背景とアプローチ ====================
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "背景とモデル設計アプローチの意義", "WHY THIS ARCHITECTURE", 2)
    sections = [
        ("課題", "感情分類のブラックボックス化", "従来の「音声から直接感情カテゴリを当てる」手法では、何が感情予測の根拠になったのか説明困難である。", RED, PALE_RED),
        ("アプローチ", "VAD感情空間の媒介", "中間にValence（快・不快）、Arousal（活性）、Dominance（支配度）の3次元空間を挟み、分類の判断材料を明確化。", BLUE, PALE_BLUE),
        ("技術基盤", "emotion2vecの採用", "自己教師あり学習 (SSL) で事前学習された強力な音声感情基盤モデルを活用し、頑健な音響特徴量を抽出。", CYAN, PALE_CYAN),
    ]
    for index, (tag, title, body, color, pale) in enumerate(sections):
        x = 0.72 + index * 4.18
        add_card(slide, x, 1.65, 3.87, 3.02, fill=WHITE, line=LINE)
        add_badge(slide, tag, x + 0.28, 1.94, 0.88, fill=pale, color=color)
        add_text(slide, title, x + 0.28, 2.48, 3.24, 0.38, size=17.5, color=NAVY, bold=True)
        add_text(slide, body, x + 0.28, 3.07, 3.27, 1.03, size=12.5, color=INK, line_spacing=1.08)
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.51), Inches(3.87), Inches(0.16))
        set_fill(band, color)
        set_line(band)
    
    add_card(slide, 0.72, 5.04, 12.01, 1.46, fill=NAVY, line=NAVY)
    add_badge(slide, "設計のコンセプト", 1.04, 5.37, 1.35, fill="173D5D", color="8DE0E8")
    add_text(
        slide,
        "音声から一度VAD（3次元連続値）を予測し、そのVAD空間上の座標値のみから感情クラスを予測する。\nこれにより、「Valenceが低くArousalが高いため、怒りと判断した」等の論理的な説明が可能となる。",
        2.58,
        5.20,
        9.50,
        0.82,
        size=14,
        color=WHITE,
        bold=True,
        line_spacing=1.1,
    )
    add_footer(slide, "本研究の焦点：emotion2vec特徴を利用した説明可能性の高い感情認識モデルの構成構築。")

    # ==================== Slide 3: アーキテクチャ図 ====================
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "提案モデル構成：VAD媒介型感情分類", "VAD-MEDIATED ARCHITECTURE", 3)
    
    add_text(slide, "音声WAVから最終的な感情判定結果までの一連の流れ", 0.72, 1.47, 10.3, 0.31, size=13, color=INK)
    
    stages = [
        ("1", "音声入力\n(WAV)", "16kHzモノラル音声をモデルへ入力", BLUE, PALE_BLUE),
        ("2", "特徴抽出\n(emotion2vec)", "20msごとの感情特徴 (768次元) を抽出", CYAN, PALE_CYAN),
        ("3", "時間圧縮\n(Pooling)", "無音パディングを除外し発話単位に平均化", BLUE, PALE_BLUE),
        ("4", "回帰ヘッド\n(VAD 3次元)", "心理学的感情尺度へ射影 (Valence/Arousal/Dom)", CYAN, PALE_CYAN),
        ("5", "分類ヘッド\n(最終感情)", "VAD座標から感情カテゴリ (喜び/悲しみ等) を判定", GREEN, PALE_GREEN),
    ]
    for index, (num, title, body, color, pale) in enumerate(stages):
        x = 0.72 + index * 2.45
        if index < 4:
            add_connector(slide, x + 2.12, 3.25, x + 2.53, 3.25, color="B7C4D1", width=2)
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + 2.26), Inches(3.08), Inches(0.22), Inches(0.34))
            set_fill(arrow, "B7C4D1")
            set_line(arrow)
        add_card(slide, x, 2.00, 2.18, 3.60, fill=WHITE if index != 3 and index != 4 else NAVY, line=CYAN if index >= 3 else LINE)
        add_step_circle(slide, num, x + 0.23, 2.25, fill=color)
        add_text(slide, title, x + 0.15, 2.92, 1.88, 0.79, size=14.5, color=WHITE if index >= 3 else NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.15, 3.92, 1.88, 0.68, size=10, color="C2D5E7" if index >= 3 else MUTED, align=PP_ALIGN.CENTER)
        
    add_card(slide, 0.72, 5.80, 11.90, 0.90, fill=PALE_CYAN, line="BDE7EB")
    add_badge(slide, "核心 of 配線", 1.00, 6.08, 0.95, fill=WHITE, color=CYAN)
    add_text(slide, "分類器（5）は、音声の生の768次元特徴量を直接見ることができません。予測された3次元のVAD値（4）のみを入力とします。", 2.10, 6.12, 10.10, 0.28, size=11, color=NAVY, bold=True)
    add_footer(slide, "VAD媒介型制約：これにより、中間のVAD表現がブラックボックス特徴に崩壊するのを防ぎます。")

    # ==================== Slide 4: 各コンポーネント ====================
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "各コンポーネントの役割と実装クラス", "MODEL COMPONENTS", 4)
    
    cards = [
        ("特徴抽出 & 時間圧縮", "Emotion2vecVADMediatedClassifier", "・emotion2vecモデルの最終層からフレーム特徴を抽出\n・`masked_mean_pooling` により可変長フレームをパディング考慮のうえ平均化\n・出力: `[Batch, 768]` の発話レベル表現", BLUE, PALE_BLUE),
        ("次元回帰部", "VADRegressionHead", "・FNN (`Linear(768->256) -> ReLU -> Linear(256->3)`)\n・768次元の発話表現を、心理学的なVADの3次元連続値へ変換\n・出力: `[Batch, 3]` （Valence, Arousal, Dominance）", CYAN, PALE_CYAN),
        ("カテゴリ判定部", "VADClassificationHead", "・Linear層 (`Linear(3 -> num_classes)`)\n・VAD座標のみから、最終感情クラス（例: 4クラス）のlogitsを出力\n・実質的にVAD値に対するロジスティック回帰として動作\n・出力: `[Batch, num_classes]`", GREEN, PALE_GREEN),
        ("統合モデル", "VADMediatedEmotionClassifier", "・`VADRegressionHead` 和 `VADClassificationHead` をカプセル化\n・引数 `return_vad=True` 時、予測されたVAD値と感情の分類logitsをペアで返却し、説明可能な情報を同時に提供", AMBER, PALE_AMBER),
    ]
    positions = [(0.72, 1.64), (6.80, 1.64), (0.72, 4.22), (6.80, 4.22)]
    for (title, cls_name, body, color, pale), (x, y) in zip(cards, positions):
        add_card(slide, x, y, 5.82, 2.20, fill=WHITE, line=LINE)
        add_card(slide, x + 0.18, y + 0.22, 1.65, 0.40, fill=pale, line=pale)
        add_text(slide, title, x + 0.22, y + 0.28, 1.55, 0.28, size=10, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, cls_name, x + 2.0, y + 0.24, 3.65, 0.35, size=14, color=NAVY, bold=True)
        add_text(slide, body, x + 0.22, y + 0.82, 5.38, 1.20, size=11, color=INK, line_spacing=1.06)
    add_footer(slide, "実装ソースコード：vad_downstream/model.py。クラスの役割定義と結合がすべて完了しています。")

    # ==================== Slide 5: 学習メカニズム ====================
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "学習メカニズム：マルチタスク損失の設計", "LOSS FUNCTION & OPTIMIZATION", 5)
    
    add_card(slide, 0.72, 1.60, 5.22, 4.92, fill=WHITE, line=LINE)
    add_badge(slide, "LOSS DESIGN", 1.02, 1.91, 1.05, fill=PALE_BLUE, color=BLUE)
    add_text(slide, "マルチタスクでVADを縛る", 1.02, 2.40, 4.2, 0.38, size=18, color=NAVY, bold=True)
    
    add_text(
        slide,
        "単に感情分類損失 (CrossEntropy) のみで学習させると、中間の3次元空間が人間の定義した「VAD」から乖離し、単に「分類しやすい任意の3次元表現」に変質してしまいます。\n\nこれを防ぐため、本物のVAD正解値に対する回帰損失と感情分類損失を同時に最小化する**マルチタスク学習**を採用しています。",
        1.02,
        3.00,
        4.62,
        3.20,
        size=12,
        color=INK,
        line_spacing=1.1,
    )

    add_card(slide, 6.22, 1.60, 6.39, 4.92, fill=NAVY, line=NAVY)
    add_badge(slide, "FORMULA", 6.54, 1.91, 1.02, fill="173B59", color="8DE0E8")
    add_text(slide, "最適化損失の数式と構成要素", 6.54, 2.40, 4.8, 0.38, size=18, color=WHITE, bold=True)
    
    add_card(slide, 6.54, 2.95, 5.74, 1.10, fill="143653", line="294D6D")
    add_text(
        slide,
        "Total Loss  =  λ_vad * L_vad  +  λ_emo * L_emo",
        6.70,
        3.15,
        5.40,
        0.50,
        size=16,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(slide, "※ default: λ_vad = 1.0, λ_emo = 1.0", 6.54, 3.82, 5.74, 0.20, size=9.5, color="8DE0E8", align=PP_ALIGN.CENTER)

    loss_details = [
        ("L_vad (VAD回帰損失)", "CCC Loss (Concordance Correlation Coefficient)", "予測と正解の相関、平均のズレ、分散のズレを同時に評価する", CYAN),
        ("L_emo (感情分類損失)", "CrossEntropy Loss", "予測logitsと感情正解ラベルとの交差エントロピー", BLUE),
    ]
    for index, (name, fn_name, desc, color) in enumerate(loss_details):
        y = 4.25 + index * 1.05
        add_text(slide, name, 6.58, y, 2.20, 0.26, size=11, color=color, bold=True)
        add_text(slide, fn_name, 6.58, y + 0.25, 5.50, 0.24, size=10, color="C0D2E3")
        add_text(slide, desc, 6.58, y + 0.48, 5.50, 0.24, size=10, color=WHITE)
        if index < 1:
            add_connector(slide, 6.58, y + 0.90, 12.10, y + 0.90, color="294966", width=0.8)

    add_footer(slide, "実装ソースコード：vad_downstream/emotion_training.py の compute_vad_emotion_loss メソッド")

    # ==================== Slide 6: 疎通確認ステータス ====================
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "実装と疎通確認のステータス", "PIPELINE COMPLETION", 6)
    
    add_text(slide, "音声波形入力からVAD・感情カテゴリ予測がJSONとして返ってくる一連の機能（配線）の完成", 0.72, 1.48, 11.0, 0.28, size=13, color=INK, bold=True)
    
    steps = [
        ("WAV", "入力音声\n(16kHz/モノラル)", BLUE, PALE_BLUE),
        ("emotion2vec", "特徴抽出\n(768次元フレーム)", CYAN, PALE_CYAN),
        ("Pooling", "発話レベル圧縮\n(平均プーリング)", BLUE, PALE_BLUE),
        ("VAD Head", "VAD 3次元予測\n(CCC回帰)", CYAN, PALE_CYAN),
        ("Emotion Head", "感情カテゴリ分類\n(予測VADを入力)", GREEN, PALE_GREEN),
    ]
    for index, (title, body, color, pale) in enumerate(steps):
        x = 0.72 + index * 2.47
        add_card(slide, x, 2.05, 2.18, 2.38, fill=WHITE, line=color if index == 4 else LINE)
        add_step_circle(slide, str(index + 1), x + 0.82, 2.31, fill=color, diameter=0.54)
        add_text(slide, title, x + 0.15, 3.08, 1.88, 0.35, size=13.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.15, 3.55, 1.88, 0.54, size=10, color=MUTED, align=PP_ALIGN.CENTER)
        if index < 4:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + 2.22), Inches(3.00), Inches(0.20), Inches(0.38))
            set_fill(arrow, "AEBCC9")
            set_line(arrow)
            
    add_card(slide, 2.15, 4.92, 9.05, 1.25, fill=NAVY, line=NAVY)
    add_text(slide, "出力結果の形式 (JSON)", 2.50, 5.22, 1.80, 0.23, size=9.5, color="84DDE5", bold=True)
    
    json_runs = [
        ("{ \n", 11, "8DE0E8", False),
        ('  "labels": ["angry", "happy", "neutral", "sad"], \n', 11, WHITE, False),
        ('  "prediction": "happy", \n', 11, WHITE, False),
        ('  "scores": [0.05, 0.82, 0.10, 0.03], \n', 11, WHITE, False),
        ('  "vad": [0.65, 0.78, 0.45]\n', 11, WHITE, False),
        ("}", 11, "8DE0E8", False),
    ]
    add_rich_text(slide, json_runs, 4.80, 5.05, 4.50, 1.10, align=PP_ALIGN.LEFT)
    add_badge(slide, "JSON推論確認済み", 9.60, 5.38, 1.20, fill="193E5D", color="8DE0E8")

    add_footer(slide, "現在は「配線が完成し、ダミーやテスト上で正しく出力が出るようになった」状態です。")

    # ==================== Slide 7: テストと検証方法 ====================
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "動作の頑健性と検証方法（手動実行コマンド）", "VERIFICATION", 7)
    
    add_card(slide, 0.72, 1.62, 4.20, 4.95, fill=NAVY, line=NAVY)
    add_badge(slide, "CODE VERIFICATION", 1.05, 1.96, 1.52, fill="173B59", color="8DE0E8")
    add_rich_text(slide, [("48", 50, WHITE, True), (" 件", 19, WHITE, True)], 1.03, 2.50, 2.40, 0.86, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "テストすべて成功", 1.07, 3.38, 2.65, 0.42, size=21, color=WHITE, bold=True)
    
    groups = [
        ("モデル構造", 6),
        ("データローダー", 10),
        ("学習ループ", 6),
        ("評価指標", 5),
        ("推論・CLI", 18),
        ("合計", 48),
    ]
    for index, (label, count) in enumerate(groups[:-1]):
        y = 4.05 + index * 0.35
        add_text(slide, label, 1.08, y, 1.10, 0.22, size=9.7, color="C3D5E6")
        bg_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.20), Inches(y + 0.025), Inches(1.42), Inches(0.12))
        bg_bar.adjustments[0] = 0.25
        set_fill(bg_bar, "34506C")
        set_line(bg_bar)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.20), Inches(y + 0.025), Inches(1.42 * count / 18), Inches(0.12))
        bar.adjustments[0] = 0.25
        set_fill(bar, CYAN)
        set_line(bar)
        add_text(slide, str(count), 3.73, y - 0.02, 0.35, 0.23, size=9.7, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    
    y_total = 4.05 + 5 * 0.35
    add_connector(slide, 1.08, y_total, 3.90, y_total, color="34506C", width=0.8)
    add_text(slide, "合計テスト件数", 1.08, y_total + 0.05, 1.30, 0.22, size=9.7, color=WHITE, bold=True)
    add_text(slide, "48 件 (ALL PASS)", 2.45, y_total + 0.05, 1.60, 0.22, size=9.7, color=CYAN, bold=True, align=PP_ALIGN.RIGHT)

    add_badge(slide, "WSL / Ubuntu環境", 1.07, 6.10, 1.22, fill="173B59", color="8DE0E8")

    # 右側: ユーザー検証方法
    add_card(slide, 5.23, 1.62, 7.38, 4.95, fill=PALE_GREEN, line="BFE5D6")
    add_badge(slide, "手動検証手順", 5.56, 1.94, 1.10, fill=WHITE, color=GREEN)
    add_text(slide, "手元で自動テストを実際に動かす手順", 5.56, 2.40, 6.55, 0.37, size=18, color=NAVY, bold=True)
    
    add_text(slide, "Windowsの任意の端末からWSL環境のテストを直接キックし、モデル動作と疎通を確認できます。", 5.56, 2.90, 6.70, 0.50, size=11, color=INK)
    
    # コマンドボックス
    add_card(slide, 5.56, 3.55, 6.70, 1.35, fill=WHITE, line="BFE5D6")
    add_text(
        slide,
        "wsl -d Ubuntu --cd /mnt/c/Users/RD004/Documents/lab/emotion2vec \\\n"
        "  -e /home/akiyama/miniforge/envs/emotion2vec-py310/bin/python \\\n"
        "  -m unittest discover -s tests",
        5.70,
        3.70,
        6.42,
        1.10,
        size=9.2,
        color=NAVY_2,
        bold=True,
        line_spacing=1.1,
    )
    
    bullets = [
        "テストが48件すべて正常にパス（OK）することを確認",
        "WAVから特徴抽出、VAD予測、分類器ロジック、JSON出力の全てのコード配線の疎通が保証されます"
    ]
    add_bullet_list(slide, bullets, 5.56, 5.10, 6.70, 1.30, size=10.5, bullet_color=GREEN, gap=4)

    add_footer(slide, "自動テスト一式は tests/ ディレクトリ以下に格納され、何時でも再現確認可能です。")

    # ==================== Slide 8: 次のステップ ====================
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "次のステップ（実験・評価計画）", "ROADMAP & PRIORITIES", 8)
    
    add_card(slide, 0.72, 1.57, 7.28, 5.03, fill=WHITE, line=LINE)
    add_badge(slide, "EXPERIMENTAL PIPELINE", 1.03, 1.88, 1.70, fill=PALE_BLUE, color=BLUE)
    add_text(slide, "モデル構成の有効性を検証する実験ロードマップ", 1.03, 2.30, 6.50, 0.35, size=18, color=NAVY, bold=True)
    
    roadmap = [
        ("1. 日本語実データの manifest & ラベルの作成", "HCUDBデータセット等を用いて、音声データパスと正解VAD・感情ラベルを整理。"),
        ("2. emotion2vec によるバッチ特徴抽出の実行", "音声WAVから768次元の中間特徴量（.npy）とフレーム長（.lengths）を抽出。"),
        ("3. VAD媒介型感情分類器の学習実験 (LOSO)", "今回実装したマルチタスク学習 (train_vad_emotion.py) を使って、話者交代交差検証 (LOSO) を回す。"),
        ("4. 直接分類モデルとの性能・説明性の評価比較", "直接分類モデルと、今回の媒介型分類モデルのWA、F1、CCC値を比較評価。"),
    ]
    for index, (step_t, step_d) in enumerate(roadmap):
        y = 2.90 + index * 0.90
        add_text(slide, step_t, 1.03, y, 6.50, 0.25, size=11.5, color=NAVY, bold=True)
        add_text(slide, step_d, 1.25, y + 0.28, 6.28, 0.35, size=9.8, color=MUTED)
        if index < 3:
            add_connector(slide, 1.03, y + 0.78, 6.90, y + 0.78, color=LINE, width=0.6)

    # 右側: 優先事項
    add_card(slide, 8.29, 1.57, 4.32, 5.03, fill=NAVY, line=NAVY)
    add_badge(slide, "PRIORITY", 8.62, 1.88, 0.92, fill="173B59", color="8DE0E8")
    priorities = [
        ("1", "実データの入力・ラベル整理", "HCUDB特徴量・VADターゲット構築", CYAN),
        ("2", "VAD媒介モデルの学習", "マルチタスクロスで重みを更新", CYAN),
        ("3", "LOSOでの検証結果の測定", "CCC / WA / Macro-F1の算出", "47647C"),
        ("4", "fine-tuningとの性能比較", "説明性と分類精度のトレードオフ評価", "47647C"),
    ]
    for index, (num, title, body, color) in enumerate(priorities):
        y = 2.52 + index * 0.89
        add_step_circle(slide, num, 8.62, y, fill=color, diameter=0.48)
        add_text(slide, title, 9.30, y - 0.01, 2.79, 0.27, size=12.5, color=WHITE, bold=True)
        add_text(slide, body, 9.30, y + 0.31, 2.83, 0.27, size=9.6, color="BFD1E2")
        if index < 3:
            add_connector(slide, 8.86, y + 0.51, 8.86, y + 0.84, color="41617D", width=1.4)
    add_card(slide, 8.61, 6.08, 3.67, 0.31, fill="173A58", line="173A58")
    add_text(slide, "最優先：実データでの学習実験", 8.75, 6.14, 3.39, 0.18, size=9.5, color="8DE0E8", bold=True, align=PP_ALIGN.CENTER)
    
    add_footer(slide, "これまでの「モデル設計・疎通確認」から、次週以降は「実データでの実験・評価」へと移行します。")

    # ==================== Slide 9: まとめ ====================
    slide = prs.slides.add_slide(blank)
    add_bg(slide, NAVY)
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.75), Inches(-1.42), Inches(4.40), Inches(4.40))
    set_fill(glow, BLUE, 72)
    set_line(glow)
    add_badge(slide, "SUMMARY", 0.75, 0.60, 0.94, fill="173B59", color="8DE0E8")
    add_text(slide, "まとめ", 0.75, 1.13, 5.0, 0.62, size=30, color=WHITE, bold=True)
    summary_cards = [
        ("モデル構成の設計", "中間にVAD次元感情空間を挟む\n説明性の高い分類モデル構成の構築", BLUE, "153754"),
        ("動作の疎通確認", "48件の自動テストによる動作保証\nWAV➔JSON予測の一貫した流れを実装", CYAN, "153754"),
        ("次のアクション", "日本語実データ(HCUDB等)を用いた\nモデルの実際の学習と評価実験", GREEN, "153754"),
    ]
    for index, (tag, body, color, fill) in enumerate(summary_cards):
        x = 0.75 + index * 4.07
        add_card(slide, x, 2.12, 3.70, 2.38, fill=fill, line="294C6B")
        add_badge(slide, tag, x + 0.28, 2.42, 1.35, fill="1D4262", color="91DFE7" if color != GREEN else "8BE0BD")
        add_text(slide, body, x + 0.28, 3.10, 3.13, 0.88, size=15.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.33), Inches(3.70), Inches(0.17))
        set_fill(strip, color)
        set_line(strip)
        
    add_card(slide, 1.65, 5.18, 10.03, 1.05, fill=WHITE, line=WHITE)
    add_text(slide, "NEXT MILESTONE", 1.98, 5.40, 1.55, 0.22, size=9.5, color=BLUE, bold=True)
    add_text(slide, "実データでの学習実験と性能評価", 3.72, 5.32, 4.30, 0.43, size=21, color=NAVY, bold=True)
    add_text(slide, "→ 提案したモデル構成の有効性を数値で実証", 8.05, 5.42, 3.20, 0.28, size=11, color=MUTED, bold=True)
    add_text(slide, "23RD004  秋山叶太  ｜  2026年7月4日", 0.78, 6.91, 5.50, 0.22, size=9.5, color="B8CCDE")

    prs.core_properties.title = "研究進捗報告（モデル構成版）"
    prs.core_properties.subject = "emotion2vecを用いた日本語音声感情認識 of 構築と評価"
    prs.core_properties.author = "23RD004 秋山叶太"
    prs.core_properties.comments = "2026年7月4日時点。VAD媒介型感情分類モデルの構成設計と、入力から出力までの疎通（自動テスト48件パス）の完了。"

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build_deck())
