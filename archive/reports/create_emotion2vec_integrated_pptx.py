from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUTPUT = ROOT / "archive" / "reports" / "emotion2vec_進捗報告_統合版.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

NAVY = "0B1F36"
NAVY_2 = "12304D"
INK = "17263A"
MUTED = "607184"
BLUE = "2C7BE5"
CYAN = "20B7C9"
GREEN = "28A878"
AMBER = "E9A23B"
RED = "D65B5B"
WHITE = "FFFFFF"
BG = "F3F6FA"
LINE = "D8E1EA"
PALE_BLUE = "EAF3FD"
PALE_CYAN = "E8F8FA"
PALE_GREEN = "EAF7F2"
PALE_AMBER = "FFF5E5"
PALE_RED = "FDEEEE"
FONT = "Yu Gothic"


def rgb(value):
    return RGBColor.from_string(value)


def set_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def set_line(shape, color=None, width=1):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(color)
        shape.line.width = Pt(width)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=16,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.02,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, runs, x, y, w, h, *, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for text, size, color, bold in runs:
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(color)


def add_card(slide, x, y, w, h, *, fill=WHITE, line=LINE, radius=True):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        card.adjustments[0] = 0.08
    set_fill(card, fill)
    set_line(card, line, 0.8)
    return card


def add_badge(slide, text, x, y, w, *, fill=PALE_BLUE, color=BLUE, h=0.32):
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    badge.adjustments[0] = 0.18
    set_fill(badge, fill)
    set_line(badge)
    add_text(
        slide,
        text,
        x,
        y + 0.005,
        w,
        h - 0.02,
        size=9,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return badge


def add_header(slide, title, kicker, number):
    add_text(slide, kicker, 0.65, 0.30, 7.0, 0.24, size=9.5, color=BLUE, bold=True)
    add_text(slide, title, 0.65, 0.62, 11.45, 0.52, size=25, color=NAVY, bold=True)
    add_text(slide, f"{number:02d}", 12.03, 0.35, 0.55, 0.28, size=11, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.22), Inches(0.70), Inches(0.045))
    set_fill(accent, CYAN)
    set_line(accent)


def add_footer(slide, text):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(7.05), Inches(12.02), Inches(0.012))
    set_fill(line, LINE)
    set_line(line)
    add_text(slide, text, 0.66, 7.11, 11.9, 0.18, size=8.2, color=MUTED)


def add_bullet_list(slide, items, x, y, w, h, *, size=13, color=INK, bullet_color=BLUE, gap=6):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.01)
    frame.margin_right = Inches(0.01)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(gap)
        paragraph.line_spacing = 1.05
        marker = paragraph.add_run()
        marker.text = "●  "
        marker.font.name = FONT
        marker.font.size = Pt(max(8, size - 4))
        marker.font.color.rgb = rgb(bullet_color)
        run = paragraph.add_run()
        run.text = item
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
    return box


def add_step_circle(slide, text, x, y, *, fill=BLUE, diameter=0.50):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(diameter), Inches(diameter))
    set_fill(circle, fill)
    set_line(circle)
    add_text(
        slide,
        text,
        x,
        y + 0.04,
        diameter,
        diameter - 0.08,
        size=11,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    return circle


def add_connector(slide, x1, y1, x2, y2, *, color=LINE, width=2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    set_line(line, color, width)
    return line


def add_status_row(slide, label, status, x, y, w, *, color, pale):
    add_card(slide, x, y, w, 0.54, fill=WHITE, line=LINE)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.17), Inches(0.20), Inches(0.20))
    set_fill(dot, color)
    set_line(dot)
    add_text(slide, label, x + 0.50, y + 0.12, w - 1.90, 0.28, size=11.5, color=INK, bold=True)
    add_badge(slide, status, x + w - 1.28, y + 0.11, 1.04, fill=pale, color=color)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    add_bg(slide, NAVY)
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.75), Inches(-1.65), Inches(5.65), Inches(5.65))
    set_fill(glow, BLUE, 68)
    set_line(glow)
    glow2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.35), Inches(4.72), Inches(3.6), Inches(3.6))
    set_fill(glow2, CYAN, 72)
    set_line(glow2)
    add_badge(slide, "RESEARCH PLAN × PROGRESS", 0.76, 0.65, 2.25, fill="173B5E", color="91DFE7")
    add_text(
        slide,
        "emotion2vecを用いた\n日本語音声感情認識の構築と評価",
        0.76,
        1.36,
        7.65,
        1.62,
        size=29,
        color=WHITE,
        bold=True,
        line_spacing=0.92,
    )
    add_text(slide, "研究計画・進捗報告 統合版", 0.80, 3.20, 5.2, 0.34, size=16, color="BFD3E7", bold=True)
    add_card(slide, 0.76, 4.12, 6.58, 1.25, fill="123250", line="284A6B")
    add_text(slide, "現在地", 1.02, 4.38, 0.86, 0.25, size=10, color="83DDE5", bold=True)
    add_text(slide, "固定特徴を用いた VA / VAD 回帰経路まで実装", 1.02, 4.73, 5.90, 0.35, size=18, color=WHITE, bold=True)
    add_text(slide, "23RD004  秋山叶太", 0.80, 6.35, 3.6, 0.31, size=14, color=WHITE, bold=True)
    add_text(slide, "2026年6月22日", 0.80, 6.76, 2.6, 0.25, size=11, color="BFD3E7")

    labels = [
        ("PLAN", "研究全体を4段階で設計", BLUE),
        ("BUILD", "VA / VAD回帰を実装", CYAN),
        ("NEXT", "実データCCCを報告", GREEN),
    ]
    for index, (tag, text, color) in enumerate(labels):
        y = 1.45 + index * 1.48
        add_card(slide, 8.75, y, 3.70, 1.10, fill=WHITE, line=WHITE)
        add_badge(slide, tag, 9.03, y + 0.20, 0.72, fill=PALE_BLUE if color == BLUE else PALE_CYAN if color == CYAN else PALE_GREEN, color=color)
        add_text(slide, text, 9.03, y + 0.61, 3.00, 0.27, size=11.5, color=NAVY, bold=True)

    # 2. Background and purpose
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "背景・先行研究・研究目的", "WHY THIS RESEARCH", 2)
    sections = [
        ("背景", "音声感情認識", "対話・教育・医療支援では、発話内容だけでなく声の感情傾向を扱う必要がある。", BLUE, PALE_BLUE),
        ("先行研究", "emotion2vec", "自己教師あり学習で汎用的な感情表現を獲得し、複数の感情認識タスクへ転用する。", CYAN, PALE_CYAN),
        ("課題", "日本語への適用", "言語・データ条件の差を踏まえ、日本語音声での連続値回帰と分類を体系的に評価する。", AMBER, PALE_AMBER),
    ]
    for index, (tag, title, body, color, pale) in enumerate(sections):
        x = 0.72 + index * 4.18
        add_card(slide, x, 1.65, 3.87, 3.02, fill=WHITE, line=LINE)
        add_badge(slide, tag, x + 0.28, 1.94, 0.88, fill=pale, color=color)
        add_text(slide, title, x + 0.28, 2.48, 3.24, 0.38, size=19, color=NAVY, bold=True)
        add_text(slide, body, x + 0.28, 3.07, 3.27, 1.03, size=13, color=INK, line_spacing=1.08)
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.51), Inches(3.87), Inches(0.16))
        set_fill(band, color)
        set_line(band)
    add_card(slide, 0.72, 5.04, 12.01, 1.46, fill=NAVY, line=NAVY)
    add_badge(slide, "研究目的", 1.04, 5.37, 1.02, fill="173D5D", color="8DE0E8")
    add_text(
        slide,
        "日本語音声に適したemotion2vec活用法を構築し、\n連続感情値回帰・感情分類・比較評価から有効性を明らかにする",
        2.38,
        5.27,
        9.70,
        0.82,
        size=19,
        color=WHITE,
        bold=True,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_footer(slide, "先行研究：emotion2vec（Ma et al., ACL 2024）。本研究は日本語条件での実装と実証を焦点とする。")

    # 3. Overall research design
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "研究全体の設計", "RESEARCH ROADMAP", 3)
    add_text(slide, "計画は4段階。現在は回帰経路の実装を終え、実データ評価へ移る地点にいる。", 0.72, 1.47, 10.3, 0.31, size=13, color=INK)
    stages = [
        ("1", "日本語\nfine-tuning", "日本語音声で\n基盤表現を適応", "未実施", AMBER, PALE_AMBER),
        ("2", "VAD回帰", "連続感情値を\nCCCで評価", "実装済み※", CYAN, PALE_CYAN),
        ("3", "感情分類", "カテゴリ感情を\nWA / Macro-F1で評価", "未実施", AMBER, PALE_AMBER),
        ("4", "比較評価", "固定特徴・fine-tuning・\n既存手法を比較", "未実施", AMBER, PALE_AMBER),
    ]
    for index, (num, title, body, status, color, pale) in enumerate(stages):
        x = 0.72 + index * 3.05
        if index < 3:
            add_connector(slide, x + 2.72, 3.25, x + 3.13, 3.25, color="B7C4D1", width=2)
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + 2.91), Inches(3.08), Inches(0.22), Inches(0.34))
            set_fill(arrow, "B7C4D1")
            set_line(arrow)
        add_card(slide, x, 2.00, 2.75, 3.52, fill=WHITE if index != 1 else NAVY, line=CYAN if index == 1 else LINE)
        add_step_circle(slide, num, x + 0.23, 2.25, fill=color)
        add_text(slide, title, x + 0.24, 2.92, 2.28, 0.79, size=18, color=WHITE if index == 1 else NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.25, 3.92, 2.24, 0.68, size=11.5, color="C2D5E7" if index == 1 else MUTED, align=PP_ALIGN.CENTER)
        add_badge(slide, status, x + 0.77, 4.86, 1.22, fill="183B59" if index == 1 else pale, color="8EE1E8" if index == 1 else color)
    add_card(slide, 3.76, 5.84, 5.80, 0.70, fill=PALE_CYAN, line="BDE7EB")
    add_text(slide, "※ 現在実装済み：固定した768次元emotion2vec特徴を使う VA / VAD 回帰", 4.00, 6.04, 5.35, 0.28, size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "emotion2vec本体のfine-tuningは未実施。現在地を研究全体の完成と混同しない。")

    # 4. Datasets and evaluation
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "データセット・評価方法", "DATA & EVALUATION", 4)
    add_card(slide, 0.72, 1.60, 5.02, 4.92, fill=WHITE, line=LINE)
    add_badge(slide, "DATA", 1.02, 1.91, 0.72, fill=PALE_BLUE, color=BLUE)
    add_text(slide, "用途でデータを分ける", 1.02, 2.40, 3.8, 0.38, size=19, color=NAVY, bold=True)
    dataset_cards = [
        ("HCUDB", "日本語音声", "日本語fine-tuning・VA/VAD回帰の候補", BLUE, PALE_BLUE),
        ("IEMOCAP", "英語・対話音声", "感情分類・既存条件との比較", CYAN, PALE_CYAN),
        ("その他公開データ", "比較用", "言語差・データ差の追加検証", GREEN, PALE_GREEN),
    ]
    for index, (name, tag, purpose, color, pale) in enumerate(dataset_cards):
        y = 3.05 + index * 0.90
        add_card(slide, 1.02, y, 4.40, 0.72, fill=pale, line=pale)
        add_text(slide, name, 1.24, y + 0.13, 1.45, 0.28, size=12, color=NAVY, bold=True)
        add_badge(slide, tag, 2.62, y + 0.18, 1.05, fill=WHITE, color=color, h=0.28)
        add_text(slide, purpose, 3.79, y + 0.13, 1.38, 0.42, size=9.2, color=INK)
    add_text(slide, "注：不一致がある発話数・分類クラス内訳は掲載しない", 1.04, 5.96, 4.22, 0.24, size=9.5, color=RED, bold=True)

    add_card(slide, 6.02, 1.60, 6.59, 4.92, fill=NAVY, line=NAVY)
    add_badge(slide, "EVALUATION", 6.34, 1.91, 1.20, fill="173B59", color="8DE0E8")
    add_text(slide, "汎化性能を分離して測る", 6.34, 2.40, 4.8, 0.38, size=19, color=WHITE, bold=True)
    add_card(slide, 6.34, 3.03, 5.94, 0.82, fill="143653", line="294D6D")
    add_text(slide, "LOSO", 6.61, 3.22, 1.02, 0.28, size=15, color="8DE0E8", bold=True)
    add_text(slide, "話者単位で学習・評価を分離", 7.72, 3.23, 3.95, 0.28, size=12.5, color=WHITE, bold=True)
    metrics = [
        ("CCC", "VA / VAD回帰", "予測と正解の一致度", CYAN),
        ("WA", "感情分類", "全体の正解率", BLUE),
        ("Macro-F1", "感情分類", "各クラスを等しく評価", GREEN),
    ]
    for index, (metric, task, desc, color) in enumerate(metrics):
        y = 4.15 + index * 0.64
        add_text(slide, metric, 6.38, y, 1.35, 0.26, size=12.5, color=color, bold=True)
        add_text(slide, task, 7.79, y + 0.01, 1.45, 0.24, size=10, color="C0D2E3")
        add_text(slide, desc, 9.46, y + 0.01, 2.55, 0.24, size=10.5, color=WHITE, bold=True)
        if index < 2:
            add_connector(slide, 6.38, y + 0.43, 12.01, y + 0.43, color="294966", width=0.8)
    add_footer(slide, "実験時はLOSOを基本とし、回帰はCCC、分類はWAとMacro-F1を主指標とする。")

    # 5. Current processing route
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "現在の処理経路", "CURRENT PIPELINE", 5)
    add_text(slide, "実装範囲", 0.72, 1.48, 1.0, 0.26, size=10, color=GREEN, bold=True)
    add_text(slide, "WAV入力から VA / VAD 出力までの配線を用意", 1.78, 1.46, 5.4, 0.28, size=13, color=INK, bold=True)
    steps = [
        ("WAV", "16 kHz\nモノラル", BLUE, PALE_BLUE),
        ("emotion2vec", "特徴抽出\n本体は固定", CYAN, PALE_CYAN),
        ("768次元特徴", "時間ごとの\nフレーム表現", BLUE, PALE_BLUE),
        ("発話単位に集約", "paddingを除外し\n平均化", CYAN, PALE_CYAN),
        ("回帰ヘッド", "小型MLP\n学習対象", GREEN, PALE_GREEN),
    ]
    for index, (title, body, color, pale) in enumerate(steps):
        x = 0.72 + index * 2.47
        add_card(slide, x, 2.05, 2.18, 2.38, fill=WHITE, line=color if index == 4 else LINE)
        add_step_circle(slide, str(index + 1), x + 0.82, 2.31, fill=color, diameter=0.54)
        add_text(slide, title, x + 0.15, 3.08, 1.88, 0.35, size=14.5, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.15, 3.55, 1.88, 0.54, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
        if index < 4:
            arrow = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + 2.22), Inches(3.00), Inches(0.20), Inches(0.38))
            set_fill(arrow, "AEBCC9")
            set_line(arrow)
    add_card(slide, 2.15, 4.92, 9.05, 1.19, fill=NAVY, line=NAVY)
    add_text(slide, "OUTPUT", 2.50, 5.22, 0.88, 0.23, size=9.5, color="84DDE5", bold=True)
    add_rich_text(
        slide,
        [
            ("VA", 22, WHITE, True),
            ("  または  ", 12, "AFC7DC", False),
            ("VAD", 22, WHITE, True),
            ("    各値 -1〜1", 12, "AFC7DC", False),
        ],
        3.65,
        5.12,
        4.95,
        0.48,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_badge(slide, "JSON推論", 9.12, 5.30, 1.14, fill="193E5D", color="8DE0E8")
    add_card(slide, 2.85, 6.34, 7.65, 0.42, fill=PALE_AMBER, line="F0D6A9")
    add_text(slide, "未実施：emotion2vec本体のfine-tuning ／ 実チェックポイントによる性能評価", 3.02, 6.45, 7.30, 0.21, size=9.8, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "現在の学習対象は回帰ヘッドのみ。emotion2vec本体は固定特徴抽出器として扱う。")

    # 6. Implemented functions
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "実装済み機能", "IMPLEMENTED", 6)
    cards = [
        ("01", "データ検証", "特徴量・長さ・ラベルの件数、次元、範囲を検査", BLUE, PALE_BLUE),
        ("02", "学習", "可変長batchからVA / VAD回帰ヘッドを学習", CYAN, PALE_CYAN),
        ("03", "最良モデル保存", "検証mean CCCが最大のcheckpointを保存", GREEN, PALE_GREEN),
        ("04", "JSON推論", "WAVと保存済みheadから予測値をJSON出力", AMBER, PALE_AMBER),
    ]
    positions = [(0.72, 1.64), (6.80, 1.64), (0.72, 4.22), (6.80, 4.22)]
    for (num, title, body, color, pale), (x, y) in zip(cards, positions):
        add_card(slide, x, y, 5.82, 2.12, fill=WHITE, line=LINE)
        add_card(slide, x + 0.22, y + 0.24, 0.76, 1.64, fill=pale, line=pale)
        add_text(slide, num, x + 0.22, y + 0.46, 0.76, 0.28, size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
        check = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.40), Inches(y + 1.17), Inches(0.40), Inches(0.40))
        set_fill(check, color)
        set_line(check)
        add_text(slide, "✓", x + 0.40, y + 1.21, 0.40, 0.26, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 1.25, y + 0.32, 3.85, 0.35, size=18, color=NAVY, bold=True)
        add_text(slide, body, x + 1.25, y + 0.88, 4.12, 0.66, size=12.5, color=INK, line_spacing=1.06)
        add_badge(slide, "実装済み", x + 4.40, y + 1.57, 1.02, fill=pale, color=color)
    add_footer(slide, "公開APIと研究コードは変更せず、既存実装の状態を資料へ反映。")

    # 7. Verification status
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "検証状況", "VERIFICATION", 7)
    add_card(slide, 0.72, 1.62, 4.20, 4.95, fill=NAVY, line=NAVY)
    add_badge(slide, "CODE VERIFICATION", 1.05, 1.96, 1.52, fill="173B59", color="8DE0E8")
    add_rich_text(slide, [("32", 50, WHITE, True), (" 件", 19, WHITE, True)], 1.03, 2.50, 2.40, 0.86, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "自動テスト成功", 1.07, 3.38, 2.65, 0.42, size=21, color=WHITE, bold=True)
    groups = [("データ", 9), ("モデル", 6), ("推論", 9), ("学習", 6), ("学習CLI", 2)]
    for index, (label, count) in enumerate(groups):
        y = 4.15 + index * 0.38
        add_text(slide, label, 1.08, y, 1.10, 0.22, size=9.7, color="C3D5E6")
        bg_bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.20), Inches(y + 0.025), Inches(1.42), Inches(0.14))
        bg_bar.adjustments[0] = 0.25
        set_fill(bg_bar, "34506C")
        set_line(bg_bar)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.20), Inches(y + 0.025), Inches(1.42 * count / 9), Inches(0.14))
        bar.adjustments[0] = 0.25
        set_fill(bar, CYAN)
        set_line(bar)
        add_text(slide, str(count), 3.73, y - 0.02, 0.35, 0.23, size=9.7, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
    add_badge(slide, "WSL / Ubuntu", 1.07, 6.05, 1.22, fill="173B59", color="8DE0E8")

    add_card(slide, 5.23, 1.62, 7.38, 2.12, fill=PALE_GREEN, line="BFE5D6")
    add_badge(slide, "確認できたこと", 5.56, 1.94, 1.30, fill=WHITE, color=GREEN)
    add_text(slide, "コードの振る舞いと入出力契約", 5.56, 2.47, 5.55, 0.37, size=19, color=NAVY, bold=True)
    add_bullet_list(slide, ["データ不整合を検出できる", "学習・評価・保存・読み込みが一連で動く"], 5.56, 2.94, 6.38, 0.66, size=11.5, bullet_color=GREEN, gap=3)

    add_card(slide, 5.23, 4.05, 7.38, 2.52, fill=PALE_AMBER, line="EFD4A7")
    add_badge(slide, "未評価", 5.56, 4.37, 0.82, fill=WHITE, color=AMBER)
    add_text(slide, "研究性能はまだ判断できない", 5.56, 4.88, 5.75, 0.38, size=19, color=NAVY, bold=True)
    add_bullet_list(
        slide,
        ["実データでのCCC・WA・Macro-F1", "実emotion2vec重みによるend-to-end疎通", "分類実験・fine-tuningの効果"],
        5.56,
        5.38,
        6.42,
        0.92,
        size=11.5,
        bullet_color=AMBER,
        gap=3,
    )
    add_footer(slide, "32件成功＝実装検証。実性能の高さを示す結果ではない。")

    # 8. Gaps and priorities
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "未完了事項・計画との差分・今後の優先順", "GAP & NEXT", 8)
    add_card(slide, 0.72, 1.57, 7.28, 5.03, fill=WHITE, line=LINE)
    add_text(slide, "計画", 1.03, 1.88, 1.15, 0.26, size=10, color=MUTED, bold=True)
    add_text(slide, "現状", 3.30, 1.88, 1.15, 0.26, size=10, color=MUTED, bold=True)
    add_text(slide, "差分", 5.26, 1.88, 1.15, 0.26, size=10, color=MUTED, bold=True)
    rows = [
        ("日本語fine-tuning", "未実施", "学習条件・重み", RED, PALE_RED),
        ("VA / VAD回帰", "head実装済み", "実データCCC", CYAN, PALE_CYAN),
        ("感情分類", "未実施", "分類器・WA / F1", RED, PALE_RED),
        ("比較評価", "未実施", "比較実験一式", RED, PALE_RED),
    ]
    for index, (plan, current, gap, color, pale) in enumerate(rows):
        y = 2.35 + index * 0.92
        add_card(slide, 1.02, y, 6.67, 0.68, fill=pale if index == 1 else BG, line=pale if index == 1 else LINE)
        add_text(slide, plan, 1.20, y + 0.19, 1.86, 0.26, size=11.5, color=NAVY, bold=True)
        add_badge(slide, current, 3.22, y + 0.18, 1.47, fill=WHITE, color=color, h=0.30)
        add_text(slide, gap, 5.10, y + 0.19, 2.26, 0.26, size=11, color=INK, bold=True)
    add_text(slide, "研究全体では『実装基盤の完成』段階。実験結果はこれから。", 1.04, 6.10, 6.43, 0.25, size=10.5, color=MUTED, bold=True)

    add_card(slide, 8.29, 1.57, 4.32, 5.03, fill=NAVY, line=NAVY)
    add_badge(slide, "PRIORITY", 8.62, 1.88, 0.92, fill="173B59", color="8DE0E8")
    priorities = [
        ("1", "実データを学習形式へ", "特徴・長さ・VADラベルを整備", CYAN),
        ("2", "回帰headを学習", "LOSOでCCCと失敗例を確認", CYAN),
        ("3", "分類経路を実装", "WA / Macro-F1を測定", "47647C"),
        ("4", "fine-tuning・比較", "固定特徴との差を検証", "47647C"),
    ]
    for index, (num, title, body, color) in enumerate(priorities):
        y = 2.52 + index * 0.89
        add_step_circle(slide, num, 8.62, y, fill=color, diameter=0.48)
        add_text(slide, title, 9.30, y - 0.01, 2.79, 0.27, size=12.5, color=WHITE, bold=True)
        add_text(slide, body, 9.30, y + 0.31, 2.83, 0.27, size=9.6, color="BFD1E2")
        if index < 3:
            add_connector(slide, 8.86, y + 0.51, 8.86, y + 0.84, color="41617D", width=1.4)
    add_card(slide, 8.61, 6.08, 3.67, 0.31, fill="173A58", line="173A58")
    add_text(slide, "最優先：実データCCC", 8.75, 6.14, 3.39, 0.18, size=9.5, color="8DE0E8", bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "優先順は『数値で現在地を示す』ことを最短距離に置く。")

    # 9. Summary
    slide = prs.slides.add_slide(blank)
    add_bg(slide, NAVY)
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.75), Inches(-1.42), Inches(4.40), Inches(4.40))
    set_fill(glow, BLUE, 72)
    set_line(glow)
    add_badge(slide, "SUMMARY", 0.75, 0.60, 0.94, fill="173B59", color="8DE0E8")
    add_text(slide, "まとめ", 0.75, 1.13, 5.0, 0.62, size=30, color=WHITE, bold=True)
    summary_cards = [
        ("達成内容", "固定emotion2vec特徴を使う\nVA / VAD回帰基盤を実装", BLUE, "153754"),
        ("現在地", "32件の自動テスト成功\n実性能は未評価", CYAN, "153754"),
        ("次の到達点", "実データCCCを\n根拠付きで報告", GREEN, "153754"),
    ]
    for index, (tag, body, color, fill) in enumerate(summary_cards):
        x = 0.75 + index * 4.07
        add_card(slide, x, 2.12, 3.70, 2.38, fill=fill, line="294C6B")
        add_badge(slide, tag, x + 0.28, 2.42, 1.15, fill="1D4262", color="91DFE7" if color != GREEN else "8BE0BD")
        add_text(slide, body, x + 0.28, 3.10, 3.13, 0.88, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(4.33), Inches(3.70), Inches(0.17))
        set_fill(strip, color)
        set_line(strip)
    add_card(slide, 1.65, 5.18, 10.03, 1.05, fill=WHITE, line=WHITE)
    add_text(slide, "NEXT MILESTONE", 1.98, 5.40, 1.55, 0.22, size=9.5, color=BLUE, bold=True)
    add_text(slide, "実データCCCの報告", 3.72, 5.32, 4.30, 0.43, size=23, color=NAVY, bold=True)
    add_text(slide, "→ 研究性能を初めて数値で判断", 8.05, 5.42, 3.08, 0.28, size=11.5, color=MUTED, bold=True)
    add_text(slide, "23RD004  秋山叶太  ｜  2026年6月22日", 0.78, 6.91, 5.50, 0.22, size=9.5, color="B8CCDE")

    prs.core_properties.title = "研究計画・進捗報告 統合版"
    prs.core_properties.subject = "emotion2vecを用いた日本語音声感情認識の構築と評価"
    prs.core_properties.author = "23RD004 秋山叶太"
    prs.core_properties.comments = "2026年6月22日時点。固定emotion2vec特徴を用いるVA/VAD回帰経路まで実装。"

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build_deck())
