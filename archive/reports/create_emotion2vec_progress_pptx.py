from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUTPUT = ROOT / "archive" / "reports" / "emotion2vec_進捗報告.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

NAVY = "10243E"
INK = "17263A"
MUTED = "5E6C7B"
BLUE = "2878D0"
CYAN = "20B7C9"
GREEN = "22A06B"
AMBER = "E99B2D"
RED = "CF4A4A"
PALE_BLUE = "EAF3FC"
PALE_CYAN = "E8F8FA"
PALE_GREEN = "EAF7F1"
PALE_AMBER = "FFF5E5"
PALE_RED = "FCECEC"
WHITE = "FFFFFF"
BG = "F5F7FA"
LINE = "D9E0E8"


def rgb(value):
    return RGBColor.from_string(value)


def set_shape_fill(shape, color, transparency=0):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def set_shape_line(shape, color=None, width=1):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(color)
        shape.line.width = Pt(width)


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    *,
    size=18,
    color=INK,
    bold=False,
    font="Yu Gothic",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.03,
    line_spacing=1.0,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, runs, x, y, w, h, *, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    for text, size, color, bold in runs:
        run = paragraph.add_run()
        run.text = text
        run.font.name = "Yu Gothic"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_bg(slide, color=BG):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(color)


def add_header(slide, title, kicker, number):
    add_text(slide, kicker, 0.65, 0.33, 6.0, 0.28, size=10, color=BLUE, bold=True)
    add_text(slide, title, 0.65, 0.65, 11.5, 0.58, size=26, color=NAVY, bold=True)
    add_text(slide, f"{number:02d}", 12.05, 0.38, 0.55, 0.33, size=12, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(1.29), Inches(0.72), Inches(0.05))
    set_shape_fill(accent, CYAN)
    set_shape_line(accent)


def add_footer(slide, text):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.65), Inches(7.08), Inches(12.03), Inches(0.012))
    set_shape_fill(line, LINE)
    set_shape_line(line)
    add_text(slide, text, 0.66, 7.13, 11.9, 0.20, size=8.5, color=MUTED)


def add_badge(slide, text, x, y, w, *, fill=PALE_GREEN, color=GREEN):
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.34))
    badge.adjustments[0] = 0.2
    set_shape_fill(badge, fill)
    set_shape_line(badge)
    add_text(slide, text, x, y + 0.01, w, 0.28, size=9, color=color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return badge


def add_card(slide, x, y, w, h, *, fill=WHITE, line=LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    card = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius:
        card.adjustments[0] = 0.08
    set_shape_fill(card, fill)
    set_shape_line(card, line, 0.8)
    return card


def add_bullet_list(slide, items, x, y, w, h, *, size=15, color=INK, bullet_color=BLUE, gap=7):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.02)
    frame.margin_right = Inches(0.02)
    frame.margin_top = Inches(0.01)
    frame.margin_bottom = Inches(0.01)
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.08
        bullet = p.add_run()
        bullet.text = "●  "
        bullet.font.name = "Yu Gothic"
        bullet.font.size = Pt(max(8, size - 4))
        bullet.font.color.rgb = rgb(bullet_color)
        run = p.add_run()
        run.text = item
        run.font.name = "Yu Gothic"
        run.font.size = Pt(size)
        run.font.color.rgb = rgb(color)
    return box


def add_number_circle(slide, number, x, y, *, fill=BLUE):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.42), Inches(0.42))
    set_shape_fill(circle, fill)
    set_shape_line(circle)
    add_text(slide, str(number), x, y + 0.005, 0.42, 0.35, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    # 1. Project purpose
    slide = prs.slides.add_slide(blank)
    add_bg(slide, NAVY)
    glow1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.35), Inches(-1.15), Inches(5.2), Inches(5.2))
    set_shape_fill(glow1, BLUE, 68)
    set_shape_line(glow1)
    glow2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(4.65), Inches(3.6), Inches(3.6))
    set_shape_fill(glow2, CYAN, 70)
    set_shape_line(glow2)
    add_badge(slide, "PROJECT PROGRESS", 0.72, 0.62, 1.75, fill="203C60", color="8EDDE5")
    add_text(slide, "emotion2vec\n進捗報告", 0.72, 1.18, 6.2, 1.65, size=34, color=WHITE, bold=True, line_spacing=0.92)
    add_text(slide, "音声から連続的な感情値を推定する", 0.76, 2.92, 6.5, 0.42, size=18, color="BFD2E7", bold=True)
    goal = add_card(slide, 0.72, 3.66, 6.35, 1.47, fill="173353", line="2B4C70")
    add_text(slide, "目的", 0.98, 3.91, 0.75, 0.30, size=11, color="78D9E4", bold=True)
    add_text(slide, "声に含まれる感情の傾向を、\n比較・学習できる数値に変える", 0.98, 4.22, 5.7, 0.72, size=22, color=WHITE, bold=True)

    labels = [
        ("Valence", "快い ↔ 不快", BLUE, PALE_BLUE),
        ("Arousal", "活発 ↔ 穏やか", CYAN, PALE_CYAN),
        ("Dominance", "強い ↔ 控えめ", AMBER, PALE_AMBER),
    ]
    for i, (name, desc, color, pale) in enumerate(labels):
        y = 1.38 + i * 1.48
        add_card(slide, 8.05, y, 4.25, 1.15, fill=WHITE, line=WHITE)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.36), Inches(y + 0.30), Inches(0.52), Inches(0.52))
        set_shape_fill(dot, color)
        set_shape_line(dot)
        add_text(slide, name, 9.10, y + 0.22, 2.6, 0.32, size=14, color=NAVY, bold=True)
        add_text(slide, desc, 9.10, y + 0.57, 2.7, 0.28, size=12, color=MUTED)
        if name == "Dominance":
            add_badge(slide, "任意", 11.48, y + 0.10, 0.55, fill=PALE_AMBER, color=AMBER)
    add_text(slide, "現在の焦点：emotion2vecの特徴量に、小さな予測器を追加する", 0.75, 6.82, 9.8, 0.26, size=10, color="AFC4DA")

    # 2. Processing mechanism
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "音声から感情値を出す仕組み", "PROCESS", 2)
    steps = [
        ("WAV音声", "16 kHz・モノラル", "音の波", BLUE),
        ("emotion2vec", "音声の特徴を抽出", "基盤モデル", CYAN),
        ("フレーム特徴", "時間ごとの768個の値", "特徴量", BLUE),
        ("発話に要約", "無音の埋め草を除外", "平均化", CYAN),
        ("予測器", "VA / VADを回帰", "小型モデル", GREEN),
    ]
    start_x = 0.72
    card_w = 2.15
    gap = 0.30
    for i, (title, body, tag, color) in enumerate(steps):
        x = start_x + i * (card_w + gap)
        add_card(slide, x, 2.05, card_w, 2.02, fill=WHITE, line=LINE)
        icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.72), Inches(2.30), Inches(0.72), Inches(0.72))
        set_shape_fill(icon, color)
        set_shape_line(icon)
        add_text(slide, str(i + 1), x + 0.72, 2.43, 0.72, 0.28, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 0.18, 3.15, card_w - 0.36, 0.34, size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.13, 3.53, card_w - 0.26, 0.40, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
        add_badge(slide, tag, x + 0.58, 4.23, 1.0, fill=PALE_BLUE if color == BLUE else (PALE_CYAN if color == CYAN else PALE_GREEN), color=color)
        if i < len(steps) - 1:
            chevron = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x + card_w + 0.04), Inches(2.79), Inches(0.22), Inches(0.44))
            set_shape_fill(chevron, "AAB6C4")
            set_shape_line(chevron)
    add_card(slide, 2.55, 4.95, 8.20, 1.22, fill=NAVY, line=NAVY)
    add_text(slide, "出力", 2.86, 5.18, 0.65, 0.28, size=11, color="81DDE6", bold=True)
    add_rich_text(
        slide,
        [
            ("Valence  ", 18, WHITE, True),
            ("・  ", 18, "7088A1", False),
            ("Arousal  ", 18, WHITE, True),
            ("・  ", 18, "7088A1", False),
            ("Dominance（任意）", 18, WHITE, True),
        ],
        3.65,
        5.17,
        6.65,
        0.36,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(slide, "各値は -1〜1 の範囲", 3.66, 5.62, 4.0, 0.24, size=10, color="BFD2E7")
    add_footer(slide, "処理経路は実装済み。実モデルの重みを使った疎通確認は未完了。")

    # 3. Implemented features
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "実装した機能", "IMPLEMENTED", 3)
    cards = [
        ("01", "データ読み込み", "特徴量・長さ・正解ラベルを読み込み、件数や範囲の不整合を検出", BLUE, PALE_BLUE),
        ("02", "感情値の予測", "長さの違う音声をまとめて扱い、VAまたはVADの値を出力", CYAN, PALE_CYAN),
        ("03", "モデルの学習・保存", "予測器のみを学習し、検証CCCが最良の重みを保存", GREEN, PALE_GREEN),
        ("04", "WAV音声からの推論", "保存済み予測器を読み込み、1つのWAVからJSONを生成", AMBER, PALE_AMBER),
    ]
    positions = [(0.72, 1.70), (6.78, 1.70), (0.72, 4.20), (6.78, 4.20)]
    for (num, title, body, color, pale), (x, y) in zip(cards, positions):
        add_card(slide, x, y, 5.82, 2.08, fill=WHITE, line=LINE)
        stripe = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.22), Inches(y + 0.25), Inches(0.64), Inches(1.58))
        stripe.adjustments[0] = 0.15
        set_shape_fill(stripe, pale)
        set_shape_line(stripe)
        add_text(slide, num, x + 0.22, y + 0.45, 0.64, 0.30, size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
        check = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.35), Inches(y + 1.14), Inches(0.38), Inches(0.38))
        set_shape_fill(check, color)
        set_shape_line(check)
        add_text(slide, "✓", x + 0.35, y + 1.18, 0.38, 0.25, size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, x + 1.15, y + 0.30, 4.2, 0.38, size=18, color=NAVY, bold=True)
        add_text(slide, body, x + 1.15, y + 0.82, 4.25, 0.76, size=13, color=INK, line_spacing=1.05)
        add_badge(slide, "実装済み", x + 4.33, y + 1.54, 1.08, fill=pale, color=color)
    add_footer(slide, "学習対象は現時点では予測器（回帰ヘッド）のみ。emotion2vec本体は固定。")

    # 4. Verification
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "検証状況", "VERIFICATION", 4)
    add_card(slide, 0.72, 1.65, 4.05, 4.83, fill=NAVY, line=NAVY)
    add_text(slide, "自動テスト", 1.08, 2.03, 2.0, 0.34, size=14, color="8EDDE5", bold=True)
    add_rich_text(slide, [("32", 52, WHITE, True), (" 件", 21, WHITE, True)], 1.04, 2.48, 2.25, 0.84, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, "すべて成功", 1.08, 3.38, 2.4, 0.44, size=22, color=WHITE, bold=True)
    groups = [("データ", 9), ("モデル", 6), ("推論", 9), ("学習", 6), ("学習CLI", 2)]
    y = 4.12
    for name, count in groups:
        add_text(slide, name, 1.08, y, 1.05, 0.22, size=10, color="C5D6E7")
        bar_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.16), Inches(y + 0.02), Inches(1.45), Inches(0.16))
        bar_bg.adjustments[0] = 0.25
        set_shape_fill(bar_bg, "36506D")
        set_shape_line(bar_bg)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.16), Inches(y + 0.02), Inches(1.45 * count / 9), Inches(0.16))
        bar.adjustments[0] = 0.25
        set_shape_fill(bar, CYAN)
        set_shape_line(bar)
        add_text(slide, str(count), 3.74, y - 0.02, 0.38, 0.24, size=10, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)
        y += 0.38
    add_badge(slide, "再実行でも成功", 1.08, 6.00, 1.56, fill="203F5D", color="91E1E8")

    add_card(slide, 5.08, 1.65, 7.53, 2.22, fill=WHITE, line=LINE)
    add_badge(slide, "評価方法", 5.42, 1.95, 1.05, fill=PALE_BLUE, color=BLUE)
    add_text(slide, "CCC（予測と正解の一致度）", 5.42, 2.43, 5.5, 0.38, size=20, color=NAVY, bold=True)
    add_text(slide, "値の動きと大きさをまとめて評価。1に近いほどよく一致します。", 5.42, 2.91, 6.45, 0.48, size=13, color=INK)
    scale_x = 5.46
    scale_y = 3.48
    scale_w = 6.15
    scale = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(scale_x), Inches(scale_y), Inches(scale_w), Inches(0.09))
    set_shape_fill(scale, LINE)
    set_shape_line(scale)
    for label, pos, color in [("-1", 0.0, RED), ("0", 0.5, AMBER), ("1", 1.0, GREEN)]:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(scale_x + scale_w * pos - 0.10), Inches(scale_y - 0.06), Inches(0.20), Inches(0.20))
        set_shape_fill(dot, color)
        set_shape_line(dot)
        add_text(slide, label, scale_x + scale_w * pos - 0.18, scale_y + 0.20, 0.36, 0.20, size=9, color=MUTED, bold=True, align=PP_ALIGN.CENTER)

    add_card(slide, 5.08, 4.17, 7.53, 2.31, fill=PALE_AMBER, line="F1D8B0")
    add_badge(slide, "現状の判定", 5.42, 4.46, 1.24, fill=WHITE, color=AMBER)
    add_text(slide, "CCCの計算・最良モデル選択は実装済み", 5.42, 4.96, 6.35, 0.38, size=18, color=NAVY, bold=True)
    add_text(slide, "実データでの性能値はまだ測定していません。\nテスト成功は、研究性能の高さを示すものではありません。", 5.42, 5.43, 6.35, 0.65, size=12.5, color=INK)
    add_footer(slide, "検証済み：コードの振る舞い。未検証：実チェックポイントと実データでの性能。")

    # 5. Current capabilities
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "現在できること", "CURRENT CAPABILITY", 5)
    add_text(slide, "2つの実行経路が利用できます", 0.72, 1.55, 5.2, 0.34, size=16, color=INK, bold=True)

    # Training route
    add_card(slide, 0.72, 2.02, 5.92, 3.92, fill=WHITE, line=LINE)
    add_badge(slide, "学習", 1.02, 2.32, 0.72, fill=PALE_GREEN, color=GREEN)
    add_text(slide, "予測器を作る", 1.02, 2.80, 3.2, 0.38, size=20, color=NAVY, bold=True)
    route1 = [
        ("特徴量 + 正解ラベル", PALE_BLUE, BLUE),
        ("予測器を学習", PALE_GREEN, GREEN),
        ("重みを保存", PALE_CYAN, CYAN),
    ]
    for i, (label, pale, color) in enumerate(route1):
        y = 3.43 + i * 0.72
        add_number_circle(slide, i + 1, 1.05, y, fill=color)
        add_card(slide, 1.62, y - 0.03, 3.96, 0.48, fill=pale, line=pale)
        add_text(slide, label, 1.82, y + 0.05, 3.55, 0.28, size=12.5, color=NAVY, bold=True)
    add_text(slide, "検証データがあれば、平均CCCが最良の重みを選択", 1.04, 5.60, 4.95, 0.25, size=10, color=MUTED)

    # Inference route
    add_card(slide, 6.78, 2.02, 5.83, 3.92, fill=WHITE, line=LINE)
    add_badge(slide, "推論", 7.08, 2.32, 0.72, fill=PALE_BLUE, color=BLUE)
    add_text(slide, "WAVから値を出す", 7.08, 2.80, 3.8, 0.38, size=20, color=NAVY, bold=True)
    route2 = [
        ("WAV + 2種類の重み", PALE_BLUE, BLUE),
        ("音声を処理", PALE_CYAN, CYAN),
        ("JSONを出力", PALE_GREEN, GREEN),
    ]
    for i, (label, pale, color) in enumerate(route2):
        y = 3.43 + i * 0.72
        add_number_circle(slide, i + 1, 7.11, y, fill=color)
        add_card(slide, 7.68, y - 0.03, 3.96, 0.48, fill=pale, line=pale)
        add_text(slide, label, 7.88, y + 0.05, 3.55, 0.28, size=12.5, color=NAVY, bold=True)
    add_text(slide, "出力：ラベル名、予測値、使用した予測器の情報", 7.10, 5.60, 4.95, 0.25, size=10, color=MUTED)

    add_card(slide, 1.75, 6.20, 9.83, 0.58, fill=PALE_AMBER, line="F0D4A5")
    add_text(slide, "注意：仮の予測器でも配線確認はできますが、その数値は研究結果として使えません。", 2.03, 6.36, 9.3, 0.27, size=11, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "実用的な推論には、emotion2vec本体の重みと学習済み予測器の両方が必要。")

    # 6. Remaining items and next work
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    add_header(slide, "未完了事項と次の作業", "NEXT STEPS", 6)
    add_card(slide, 0.72, 1.65, 5.10, 4.93, fill=WHITE, line=LINE)
    add_badge(slide, "未完了", 1.02, 1.96, 0.92, fill=PALE_RED, color=RED)
    add_bullet_list(
        slide,
        [
            "実emotion2vec重みを使った疎通確認",
            "実データでのCCC測定と比較",
            "生の注釈から学習形式への変換",
            "音声の再サンプリング・ステレオ対応",
            "emotion2vec本体の追加学習",
        ],
        1.02,
        2.54,
        4.40,
        3.35,
        size=13.5,
        bullet_color=RED,
        gap=10,
    )
    add_text(slide, "分類機能や大規模な実験管理も将来候補", 1.03, 6.06, 4.35, 0.25, size=9.5, color=MUTED)

    add_card(slide, 6.12, 1.65, 6.49, 4.93, fill=NAVY, line=NAVY)
    add_badge(slide, "優先順", 6.48, 1.96, 0.92, fill="203F5D", color="91E1E8")
    next_steps = [
        ("1", "必要な重みを準備", "本体 + 学習済み予測器"),
        ("2", "実WAVで疎通確認", "入力からJSONまで通す"),
        ("3", "実データでCCC評価", "性能値と失敗例を確認"),
        ("4", "前処理と対応形式を拡張", "実運用に必要な堅牢性を追加"),
    ]
    y = 2.55
    for num, title, body in next_steps:
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.50), Inches(y), Inches(0.55), Inches(0.55))
        set_shape_fill(circle, CYAN if num in ("1", "2") else "3C5875")
        set_shape_line(circle)
        add_text(slide, num, 6.50, y + 0.08, 0.55, 0.30, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, 7.30, y - 0.01, 4.65, 0.32, size=15, color=WHITE, bold=True)
        add_text(slide, body, 7.30, y + 0.32, 4.65, 0.27, size=10.5, color="BED0E2")
        if num != "4":
            stem = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.76), Inches(y + 0.57), Inches(0.035), Inches(0.42))
            set_shape_fill(stem, "4A6785")
            set_shape_line(stem)
        y += 0.93
    add_card(slide, 6.47, 6.03, 5.73, 0.34, fill="183653", line="183653")
    add_text(slide, "次の到達点：実データのCCCを根拠付きで報告する", 6.66, 6.08, 5.34, 0.20, size=10.5, color="8FE2E9", bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, "実装済みと未完了を分離し、次は実環境・実データでの検証へ進む。")

    prs.core_properties.title = "emotion2vec 進捗報告"
    prs.core_properties.subject = "目的・実装内容・検証状況・次の作業"
    prs.core_properties.author = "emotion2vec project"
    prs.core_properties.comments = "Git履歴と実装・テスト結果を確認して作成"

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    print(build_deck())
