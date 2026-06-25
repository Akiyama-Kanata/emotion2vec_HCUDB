from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[2]
PPTX = ROOT / "archive" / "reports" / "emotion2vec_進捗報告_統合版.pptx"

EXPECTED_TITLES = [
    "emotion2vecを用いた",
    "背景・先行研究・研究目的",
    "研究全体の設計",
    "データセット・評価方法",
    "現在の処理経路",
    "実装済み機能",
    "検証状況",
    "未完了事項・計画との差分・今後の優先順",
    "まとめ",
]

REQUIRED_BY_SLIDE = [
    ["23RD004", "秋山叶太", "2026年6月22日"],
    ["背景", "先行研究", "研究目的", "Ma et al."],
    ["日本語", "fine-tuning", "VAD回帰", "感情分類", "比較評価", "未実施"],
    ["HCUDB", "IEMOCAP", "LOSO", "CCC", "WA", "Macro-F1"],
    ["WAV", "emotion2vec", "768次元特徴", "回帰ヘッド", "VA", "VAD", "未実施"],
    ["データ検証", "学習", "最良モデル保存", "JSON推論"],
    ["32", "自動テスト成功", "研究性能", "未評価"],
    ["計画との差分", "実データCCC"],
    ["達成内容", "現在地", "次の到達点", "実データCCC"],
]


def slide_text(slide):
    return "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text_frame"))


def verify():
    prs = Presentation(PPTX)
    errors = []

    if len(prs.slides) != 9:
        errors.append(f"slide count: expected 9, got {len(prs.slides)}")

    expected_width = 13.333
    expected_height = 7.5
    width = prs.slide_width / 914400
    height = prs.slide_height / 914400
    if abs(width - expected_width) > 0.002 or abs(height - expected_height) > 0.002:
        errors.append(f"slide size: expected 13.333 x 7.5, got {width:.3f} x {height:.3f}")

    for index, slide in enumerate(prs.slides):
        text = slide_text(slide)
        if EXPECTED_TITLES[index] not in text:
            errors.append(f"slide {index + 1}: missing title {EXPECTED_TITLES[index]!r}")
        for required in REQUIRED_BY_SLIDE[index]:
            if required not in text:
                errors.append(f"slide {index + 1}: missing required text {required!r}")

        for shape in slide.shapes:
            # Decorative circles deliberately bleed past the canvas. Text content may not.
            if not hasattr(shape, "text_frame") or not shape.text.strip():
                continue
            if shape.left < 0 or shape.top < 0:
                errors.append(f"slide {index + 1}: text starts outside slide")
            if shape.left + shape.width > prs.slide_width + 1000:
                errors.append(f"slide {index + 1}: text exceeds right edge")
            if shape.top + shape.height > prs.slide_height + 1000:
                errors.append(f"slide {index + 1}: text exceeds bottom edge")
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.text.strip() and run.font.name not in (None, "Yu Gothic"):
                        errors.append(
                            f"slide {index + 1}: unexpected font {run.font.name!r} in {run.text!r}"
                        )

    all_text = "\n".join(slide_text(slide) for slide in prs.slides)
    if "fine-tuningは未実施" not in all_text:
        errors.append("fine-tuning must be explicitly marked as not performed")
    if "実性能は未評価" not in all_text and "実性能の高さを示す結果ではない" not in all_text:
        errors.append("real-data performance must be explicitly marked as unevaluated")

    if errors:
        print("FAILED")
        for error in errors:
            print(f"- {error}")
        raise SystemExit(1)

    print("OK")
    print(f"slides={len(prs.slides)}")
    print(f"size={width:.3f}x{height:.3f}")
    print("date=2026年6月22日")
    print("key_metrics=LOSO, CCC, WA, Macro-F1")
    print("key_numbers=768, 32")


if __name__ == "__main__":
    verify()
